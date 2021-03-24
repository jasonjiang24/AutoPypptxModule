# -*- coding: utf-8 -*-

from pptx.chart.data import CategoryChartData
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor


class EleUtils:
    def __init__(self):
        pass
    def rgb_color(self,color_list):
        return RGBColor(color_list[0],color_list[1],color_list[2])
    def fill_set(self,fill_obj,fill_pars = {}):
        '''
        #设置背景填充
        #参数示例
        fill_pars = {'transparent':False, #是否透明
                     'color':[0,0,0]} #颜色
        '''
        if bool(fill_pars):
            if fill_pars.get('transparent'):
                fill_obj.fill.background()
            else:
                fill_obj.fill.solid()
                fill_obj.fill.fore_color.rgb = EleUtils.rgb_color(self,fill_pars.get('color'))
        else:
            pass
    def line_set(self,line_obj,line_pars = {}):
        '''
        #设置线框
        #参数示例：
        line_pars = {'type':'', #设置线的类型
                     'width':4, #设置线的宽度
                     'fill':{}, #设置线的颜色，调用fill_set的设置
                    }
        '''
        if bool(line_pars):
            if bool('type'):
                line_obj.line.dash_style = line_pars.get('type')
            if bool('width'):
                line_obj.line.width = Pt(line_pars.get('width'))
            if bool('fill'):
                EleUtils.fill_set(self,line_obj.line,line_pars.get('fill'))
        else:
            pass
    def text_set(self,text_obj,text_pars = {}):
        '''
        #对TEXT文本的样式设计
        #参数示例：
        text_pars = {'alignment':'', #对其设置
                     'size':10, #大小，内置Pt()
                     'bold':False, #粗体
                     'color':[0,0,0], #颜色
                     'italic':False, #斜体
                     'language':1, #字体
                     'content':'',文本内容
                     }
        '''
        if bool(text_pars):
            if bool(text_pars.get('content')):
                text_obj.text = text_pars.get("content")
            if bool(text_pars.get('alignment')):
                text_obj.alignment = text_pars.get("alignment")
            if bool(text_pars.get('size')):
                text_obj.font.size = Pt(text_pars.get("size"))
            if bool(text_pars.get('bold')):
                text_obj.font.bold = text_pars.get("bold")
            if bool(text_pars.get('color')):
                text_obj.font.color.rgb = EleUtils.rgb_color(self,text_pars.get("color"))                
            if bool(text_pars.get('italic')):
                text_obj.font.italic = text_pars.get("italic")
            if bool(text_pars.get('language')):
                text_obj.font.language_id = text_pars.get("language")                               
        else:
            pass
    def chart_title_set(self,chart,title_pars={}):
        '''
        #设置图的名称
        #参数示例：
        title_pars = {'text':'', #图表名称
                      'text_set':{},#图表名称样式属性
                      }
        '''
        if bool(title_pars):
            title = chart.chart_title.text_frame
            title.text = title_pars.get('text')
            if bool(title_pars.get('text_set')):
                EleUtils.text_set(self,
                                  title.paragraphs[0],
                                  title_pars.get('text_set'))
        else:
            pass         
    def chart_label_set(self,label_obj,label_pars={}):
        '''
        #设置图表的数据标签
        #参数示例：
        label_pars = {'show':True, #是否显示数据标签
                      'text':{}, #字体格式，调用text_set,无需输入文本内容
                      'format':'', #数字显示格式，尤其是%显示
                      'position':'', #数据标签显示位置
                          }
        '''
        if bool(label_pars):
            if bool(label_pars.get('show')):
                label_obj.show_value = True
                if bool(label_pars.get('text')):
                    EleUtils.text_set(self,label_obj,label_pars.get('text'))
                if bool(label_pars.get('format')):
                    label_obj.number_format = label_pars.get('format')
                if bool(label_pars.get('position')):
                    label_obj.position = label_pars.get('position')
        else:
            pass           
    def chart_series_set(self,chart,series_pars = {}):
        '''
        #设置数据系列格式，包括柱颜色，线颜色，数据标签
        #参数示例：
        series_pars = {'fill_set':[{},{}], #设置系列的背景属性,调用fill_set
                       'line_set':[{},{}], #设置系列的线条属性,调用line_set
                       'label_set':[{},{}]}, #设置系列的标签属性,chart_label_set
        
        '''
        if bool(series_pars):
            if bool(series_pars.get('fill_set')):
                for fill_num in range(len(series_pars.get('fill_set'))):
                    if bool(series_pars.get('fill_set')[fill_num]):
                        EleUtils.fill_set(self,
                                          fill_obj=chart.series[fill_num].format,
                                          fill_pars=series_pars.get('fill_set')[fill_num])
            if bool(series_pars.get('line_set')):
                for line_num in range(len(series_pars.get('line_set'))):
                    if bool(series_pars.get('line_set')[line_num]):
                        EleUtils.line_set(self,
                                          line_obj=chart.series[line_num].format,
                                          line_pars=series_pars.get('line_set')[line_num])                    
            if bool(series_pars.get('label_set')):
                for label_num in range(len(series_pars.get('label_set'))):
                    if bool(series_pars.get('label_set')[label_num]):
                        EleUtils.chart_label_set(self,
                                                 label_obj=chart.series[label_num].data_labels,
                                                 label_pars=series_pars.get('label_set')[label_num])
        else:
            pass
    def chart_axis_set(self,axis_obj,axis_pars = {}):
        '''
        #设置图的X坐标轴
        #参数示例：
        xaxis_pars = {'text_show':True, #是否显示轴标签
                      'text_format':{}, #轴标签字体格式
                      'name_show':True, #是否显示轴标题
                      'name_format':{}, #轴标题字体格式
                      }
        '''
        if bool(axis_pars):
            if bool(axis_pars.get('text_show')):
                EleUtils.text_set(self,
                                  axis_obj.tick_labels,
                                  axis_pars.get('text_format'))
            else:
                axis_obj.visible=False
            if bool(axis_pars.get('name_show')):
                EleUtils.text_set(self,
                                  axis_obj.axis_title.text_frame.paragraphs[0],
                                  axis_pars.get('name_format'))
        else:
            pass
    def chart_legend_set(self,chart_obj,legend_pars = {}):
        '''
        #设置图例
        #参数示例：
        ledend_pars = {'show':True, #是否显示图例
                       'position':'', #图例位置
                       'text':{}, #图例文本样式
                       } 
        '''
        if bool(legend_pars):
            if bool(legend_pars.get('show')):
                chart_obj.has_legend = True
                if bool(legend_pars.get('position')):
                    chart_obj.legend.position = legend_pars.get('position')
                if bool(legend_pars.get('text')):
                    EleUtils.text_set(self,
                                      chart_obj.legend,
                                      legend_pars.get('text'))
            else:
                chart_obj.has_legend = False
        else:
            pass
class ChartElements:
    def __init__(self):
        self.utils = EleUtils()
    def general_chart(self,slide,chart_pars = {}):
        '''
        #创建通用图元素，如线性图，柱状图，饼图等.
        #参数示例：
        chart_pars = {'title':{}, #图表名称
                      'type':'bar', #图表类型
                      'position':[a,b,c,d], #位置4个参数，分别为纵坐标，横坐标，宽，高
                      'data':{'df':df, #数据源，为DataFrame对象,表头只能是单行,不能多行结构
                              'x_name':'',#X轴数据列
                              'series_names':['a','b','c'],#Y周系列的数据列
                              },
                      'series':{}, #设置系列的标签属性,调用chart_series_set
                      'x_axis':{}, #设置X轴
                      'y_axis':{}, #设置Y轴
                      'legend':{}, #设置图例
                      }
        '''
        chart_data = CategoryChartData()
        #设置数据
        df = chart_pars.get('data').get('df')
        #X轴数据
        chart_data.categories = df[chart_pars.get('data').get('x_name')]
        #Y轴数据
        for i in chart_pars.get('data').get('series_names'):
            chart_data.add_series(i,df[i])  
        #创建图
        graphic_frame = slide.shapes.add_chart(chart_pars.get('type'), 
                                               Inches(chart_pars.get('position')[0]), 
                                               Inches(chart_pars.get('position')[1]), 
                                               Inches(chart_pars.get('position')[2]), 
                                               Inches(chart_pars.get('position')[3]), 
                                               chart_data)
        chart = graphic_frame.chart
        #设置图表名称
        if bool(chart_pars.get('title')):
            self.utils.chart_title_set(chart,chart_pars.get('title'))
        
        #设置系列的样式
        if bool(chart_pars.get('series')):
            self.utils.chart_series_set(chart,chart_pars.get('series'))
            
        #设置x轴
        if bool(chart_pars.get('x_axis')):
            self.utils.chart_axis_set(chart.category_axis,chart_pars.get('x_axis'))
            
        #设置y轴
        if bool(chart_pars.get('y_axis')):
            self.utils.chart_axis_set(chart.value_axis,chart_pars.get('y_axis'))
        
        #设置图例
        if bool(chart_pars.get('legend')):
            self.utils.chart_legend_set(chart,chart_pars.get('legend'))
        
        return chart
    def sandiantu(self):
        pass

class ShapeElements:
    def __init__(self):
        self.utils = EleUtils()
    def general_shape(self,slide,shape_pars = {}):
        '''
        #创建通用的shape
        #参数示例
        shape_pars = {'type':'', #shape类型
                      'position':[1,2,3,4],#位置及大小
                      'fill':{},#设置背景色
                      'line':{},#设置线框
                      'text':{},#设置文本内容
                      }
        '''
        shapes = slide.shapes
        shape = shapes.add_shape(shape_pars.get('type'),
                                Inches(shape_pars.get('position')[0]),
                                Inches(shape_pars.get('position')[1]),
                                Inches(shape_pars.get('position')[2]),
                                Inches(shape_pars.get('position')[3]))
        #设置背景
        if bool(shape_pars.get('fill')):
            self.utils.fill_set(shape,shape_pars.get('fill'))
        #设置线框
        if bool(shape_pars.get('line')):
            self.utils.line_set(shape,shape_pars.get('line'))
        #设置文本内容
        if bool(shape_pars.get('text')):
            self.utils.text_set(shape.text_frame.paragraphs[0],
                                shape_pars.get('text'))
        return shape

class TableElements:
    def __init__(self):
        self.utils= EleUtils()
    def general_table(self,slide,table_pars = {}):
        '''
        #创建通用表格
        #参数示例：
        table_pars = {'position':[a,b,c,d],#表格位置&单元格大小
                      'fill_color':[[0,0,0],[0,0,0],[0,0,0]], #表格颜色，分别为表头，奇数行，偶数行颜色
                      'data':df, #DataFrame,只支持单行表头的表
                      'text':{},#单元格字体样式,调用 text_set
                      }
        '''
        if bool(table_pars):
            from pptx.enum.text import MSO_AUTO_SIZE
            df = table_pars.get('data')
            row,col = df.shape
            shapes = slide.shapes
            shape = shapes.add_table(row+1,
                                     col,
                                     Inches(table_pars.get('position')[0]),
                                     Inches(table_pars.get('position')[1]),
                                     Inches(table_pars.get('position')[2]),
                                     Inches(table_pars.get('position')[3]),
                                     )
            table = shape.table
            head_start = 0
            #表头
            for head_cell_num in range(col):
                cell = table.cell(head_start,head_cell_num)
                cell.text = df.columns.to_list()[head_cell_num]
                cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.utils.rgb_color(table_pars.get('fill_color')[0])
                if bool(table_pars.get('text')):
                    self.utils.text_set(cell.text_frame.paragraphs[0],table_pars.get('text'))
            #表内容
            for da_i in range(row):
                for da_j in range(col):
                    cell = table.cell(1+da_i,da_j)
                    cell.text = str(df.values[da_i][da_j])
                    cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    if bool(table_pars.get('text')):
                        self.utils.text_set(cell.text_frame.paragraphs[0],table_pars.get('text'))
                    if da_i%2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.utils.rgb_color(table_pars.get('fill_color')[1])
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.utils.rgb_color(table_pars.get('fill_color')[2])                        
            return shape
        
        
        
        
        
        
        
        
        
        
        
        
        
        
        