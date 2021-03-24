# -*- coding: utf-8 -*-

from pptx import Presentation
from ElementsModules import ChartElements,ShapeElements,TableElements


class ElementsRoute:
    def __init__(self):
        self.chart = ChartElements()
        self.shape = ShapeElements()
        self.table = TableElements()
    def elements_route(self,slide,elements_pars = {}):
        '''
        #PPT页内的元素路由，根据元素的类型，选取不同的函数或类
        #参数示例：
        elements_pars = {'type':'table',
                         'parameters':None}
        #type指元素的类型，常用为3大类chart,table,shape(shape可设置为图形或文本框)
        #parameters指元素对应的参数内容
        '''
        if bool(elements_pars):
            if elements_pars.get('type') == 'chart':
                self.chart.general_chart(slide,elements_pars.get('parameters'))
                return slide
            elif elements_pars.get('type') == 'table':
                self.table.general_table(slide,elements_pars.get('parameters'))
                return slide
            elif elements_pars.get('type') == 'shape':
                self.shape.general_shape(slide,elements_pars.get('parameters'))
                return slide
        else:
            return slide

class SlidesGenerate:
    def __init__(self):
        self.elements = ElementsRoute()
    def slidesgenerate(self,presentation,slides_pars = {}):
        '''
        #创建ppt内部的页
        #参数示例：
        slides_pars = {'slide1':{'back_type':0,'elements':None},
                       'slide2':{'back_type':0,'elements':None},
                       'slide3':{'back_type':0,'elements':None},
                       }
        #slides1~3指PPT的哪一页，从1~N的递增顺序增加页数
        #back_type指该页PPT用的是哪个默认背景
        #elements指该页ppt内的元素的参数
        '''
        if bool(slides_pars):
            for i in slides_pars.keys():
                background = slides_pars.get(i).get('back_type')
                slide = presentation.slides.add_slide(presentation.slide_layouts[background])
                if slides_pars.get(i).get('elements'):
                    self.elements.elements_route(slide=slide,
                                                 elements_pars=slides_pars.get(i).get('elements'))
                    pass           
            return presentation   
        else:
            return presentation

class PresentationGenerate:
    def __init__(self):
        self.slides_generate = SlidesGenerate()
    def ppt_generate(self,present_pars = {}):
        '''
        #创建PPT并保存。
        #参数示例：
        present_pars = {'input_template':None,
                        'output_save':'output_ppt.pptx',
                        'slides_info':None}
        #input_template默认为None,若有值，则为输入PPT模板的路径名称
        #output_save为必须，为输出生成的ppt路径名称
        #slides_info为内部PPT页的参数内容，若为None，则为空PPT,或复制PPT模板
        '''
        prs = Presentation(present_pars.get('input_template',None))
        
        if present_pars.get('slides_info',None) != None:
            #创建ppt页内容
            prs = self.slides_generate.slidesgenerate(presentation = prs,
                                                      slides_pars = present_pars.get('slides_info')
                                                      )

        prs.save(present_pars.get('output_save'))
        return print('ppt generated')



from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN

shape_pars =  {'type':MSO_SHAPE.RECTANGLE,
               'position':[1,1,1,1],
               'fill':{'transparent':False,
                       'color':[111,111,111]},
               'line':{'type':MSO_LINE.SOLID, 
                        'width':1, 
                        'fill':{'transparent':False,
                                 'color':[0,0,0]}, 
                        },
               'text':{'content':'aserdghgfkdjfd',
                         'alignment':PP_ALIGN.CENTER, 
                         'size':28, 
                         'bold':True, 
                         'color':[0,0,0], 
                         'italic':True, 
                         },
               }

import pandas as pd
import numpy as np
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION



dates = pd.date_range("20130101", periods=6)
df = pd.DataFrame(np.random.randint(5, 100,(6,4)), index=dates, columns=list("ABCD")).reset_index()

chart_pars = {'title':{'text':'aaaaaaaaaaaaaaaa', 
                      'text_set':{},
                      }, 
              'type':XL_CHART_TYPE.LINE, 
              'position':[1,1,6,6], 
              'data':{'df':df, 
                      'x_name':'index',
                      'series_names':['A','B','D']},
              'series':{'fill_set':[{'transparent':False,
                                     'color':[55,55,55]},
                                    {'transparent':False,
                                     'color':[222,222,222]}], 
                        'line_set':[{'width':8},{'width':4}], 
                        'label_set':[{'show':True},
                                     {'show':True}]
                        }, 
              'x_axis':{'text_show':False,
                      'text_format':{}, 
                      'name_show':True, 
                      'name_format':{'content':'bbbbuibb'}, }, 
              'y_axis':{'text_show':False,
                      'text_format':{}, 
                      'name_show':True, 
                      'name_format':{'content':'cccccc'}, }, 
              'legend':{'show':False, 
                       'position':XL_LEGEND_POSITION.TOP, 
                       'text':{}, 
                       } , 
              }

table_pars = {'position':[1,1,6,1],
              'fill_color':[[0,139,139],[121,205,205],[0,205,205]], 
              'data':df, 
              'text':{'size':8},
                      }

elements_pars = {'type':'chart',
                 'parameters':chart_pars}

slides_pars = {'slide1':{'back_type':0,'elements':elements_pars},
               'slide2':{'back_type':1,'elements':elements_pars},
               'slide3':{'back_type':2,'elements':None},
               }

present_pars = {'input_template':None,
                'output_save':'output_ppt.pptx',
                'slides_info':slides_pars}

a = PresentationGenerate()
b = a.ppt_generate(present_pars)
